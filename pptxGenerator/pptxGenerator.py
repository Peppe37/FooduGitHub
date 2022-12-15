from pptx import Presentation

def titleSlide(mytitle,mysubtitle):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
        
    title.text = mytitle
    subtitle.text = mysubtitle


prs = Presentation()

titleSlide('Report di Ricerca','Minardo - Pastina Baby Food')

prs.save('test.pptx')